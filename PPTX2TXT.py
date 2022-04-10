from pptx import Presentation

prs = Presentation('52.pptx')
for slide in prs.slides :
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)