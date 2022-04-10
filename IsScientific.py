import os
from docx import Document
import PyPDF2
from pptx import Presentation


def docx_to_dict(file_name):
  docx_dict = {}
  document = Document(file_name)
  indx = 0
  for para in document.paragraphs:
    indx += 1
    if (len(para.text) > 0):
      docx_dict[indx] = para.text
  return docx_dict

def texting(doc):
    print(doc)
    df = docx_to_dict(doc)
    txt = ""
    for pp in df.values():
      txt += pp + "\n"
    return txt

def pdfinf(chaine):
  print(chaine)
  tt = ""
  pdf = PyPDF2.PdfFileReader(chaine)
  for i in range(pdf.numPages):
    page_1_object = pdf.getPage(i)
    page_1_text = page_1_object.extractText()
    tt += page_1_text
  return tt
def pptxing (chaine):
  print(chaine)
  tt = ""
  prs = Presentation(chaine)
  for slide in prs.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        tt += shape.text
  return tt


scientifcTxt = ""
nonScientifcTxt = ""

scientificPath = "./Scientific"
NonScientificPath = "./NonScientific"

nonScientifcFiles = os.listdir(NonScientificPath)
ScientifcFiles = os.listdir(scientificPath)

for file in nonScientifcFiles:
  if file.endswith(".pdf"):
    nonScientifcTxt += pdfinf(NonScientificPath+"/"+str(file))
  elif file.endswith(".pptx"):
    nonScientifcTxt += pptxing(NonScientificPath+"/"+str(file))
  elif file.endswith(".docx"):
    nonScientifcTxt += texting(NonScientificPath+"/"+str(file))

for file in ScientifcFiles:
  if file.endswith(".pdf"):
    scientifcTxt += pdfinf(scientificPath+"/"+str(file))
  elif file.endswith(".pptx"):
    scientifcTxt += pptxing(scientificPath+"/"+str(file))
  elif file.endswith(".docx"):
    scientifcTxt += texting(scientificPath+"/"+str(file))


with open('nonScientific.txt', 'w') as f:
    f.writelines(nonScientifcTxt)

with open('scientific.txt', 'w') as f:
    f.writelines(scientifcTxt)